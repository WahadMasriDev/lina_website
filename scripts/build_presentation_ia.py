# -*- coding: utf-8 -*-
"""Génère 'presentation_IA.pptx' — « De ChatGPT au Multi-Agent ».
Présentation éducative (élèves de seconde), 16:9, schémas en formes natives.
Aucune image, aucun emoji. Voir docs/superpowers/specs/2026-06-21-...-design.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------
# PALETTE & GÉOMÉTRIE
# --------------------------------------------------------------------------
CREAM = RGBColor(0xF7, 0xF5, 0xEF)   # fond
PAPER = RGBColor(0xFF, 0xFF, 0xFF)   # cartes claires
INK   = RGBColor(0x23, 0x23, 0x2A)   # texte principal
MUTE  = RGBColor(0x6A, 0x6A, 0x74)   # texte secondaire
SLATE = RGBColor(0x2F, 0x4A, 0x6B)   # accent structurant
HAIR  = RGBColor(0xD9, 0xD5, 0xCA)   # filets / bordures discrètes

# Code couleur des rôles (réutilisé sur tout le deck)
PLAN = RGBColor(0x2F, 0x6B, 0xB0)    # Planificateur — bleu
DEV  = RGBColor(0x6B, 0x4F, 0xA0)    # Développeur   — violet
TEST = RGBColor(0x2E, 0x8B, 0x57)    # Testeur       — vert
FIX  = RGBColor(0xC9, 0x7B, 0x1E)    # Correcteur    — orange

FONT = "Segoe UI"

PW, PH = 13.333, 7.5                  # pouces, 16:9
ML, MR = 0.75, 0.75                   # marges latérales
CW = PW - ML - MR                     # largeur de contenu


def tint(rgb, p):
    """Éclaircit une couleur vers le blanc (p entre 0 et 1)."""
    return RGBColor(
        int(rgb[0] + (255 - rgb[0]) * p),
        int(rgb[1] + (255 - rgb[1]) * p),
        int(rgb[2] + (255 - rgb[2]) * p),
    )


# --------------------------------------------------------------------------
# HELPERS BAS NIVEAU
# --------------------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(PW)
    prs.slide_height = Inches(PH)
    return prs


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # disposition vide
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CREAM
    return s


def _no_shadow(shape_):
    shape_.shadow.inherit = False


def para(tf, runs, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
         first=False, space_after=4, space_before=0, line=1.04, font=FONT):
    """Ajoute un paragraphe. `runs` = str | [(txt,bold)] | [(txt,bold,color)]."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for item in runs:
        if len(item) == 2:
            txt, b = item
            col = color
        else:
            txt, b, col = item
        r = p.add_run()
        r.text = txt
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = col
    return p


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def shape(s, kind, x, y, w, h, fill=None, line_color=None, line_w=1.0,
          radius=None, anchor=MSO_ANCHOR.MIDDLE):
    sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sh)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return sh


def card(s, x, y, w, h, fill=PAPER, line_color=HAIR, line_w=1.25, radius=0.045):
    return shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                 fill=fill, line_color=line_color, line_w=line_w, radius=radius)


def rect(s, x, y, w, h, fill, line_color=None, line_w=1.0):
    return shape(s, MSO_SHAPE.RECTANGLE, x, y, w, h,
                 fill=fill, line_color=line_color, line_w=line_w)


def line(s, x1, y1, x2, y2, color=SLATE, w=2.0, arrow=True, dash=None):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    ln.set('cap', 'rnd')
    if dash is not None:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    if arrow:
        end = ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'h': 'med'})
        ln.append(end)
    _no_shadow(conn)
    return conn


def check_badge(s, cx, cy, r=0.17, color=TEST):
    """Pastille ronde avec une coche dessinée (formes, pas d'emoji)."""
    disc = shape(s, MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r, fill=color)
    line(s, cx - 0.085, cy + 0.005, cx - 0.02, cy + 0.075,
         color=PAPER, w=2.2, arrow=False)
    line(s, cx - 0.02, cy + 0.075, cx + 0.10, cy - 0.075,
         color=PAPER, w=2.2, arrow=False)
    return disc


# --------------------------------------------------------------------------
# GABARITS
# --------------------------------------------------------------------------
def footer(s, n):
    line(s, ML, 7.06, PW - MR, 7.06, color=HAIR, w=1.0, arrow=False)
    rect(s, ML, 7.05, 0.34, 0.03, fill=SLATE)
    _, tf = textbox(s, PW - MR - 3.0, 7.12, 3.0, 0.3)
    para(tf, [("De ChatGPT au Multi-Agent", False, MUTE),
              ("    %02d / 14" % n, True, SLATE)],
         size=9, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def header(s, kicker, title, kcolor=SLATE):
    rect(s, 0, 0, 0.22, PH, fill=SLATE)              # barre latérale
    _, tf = textbox(s, ML, 0.52, CW, 0.34)
    para(tf, kicker.upper(), size=12.5, bold=True, color=kcolor,
         first=True, space_after=0)
    _, tf = textbox(s, ML, 0.86, CW, 0.92)
    para(tf, title, size=30, bold=True, color=INK,
         first=True, space_after=0, line=1.0)
    line(s, ML, 1.64, ML + 0.9, 1.64, color=kcolor, w=2.6, arrow=False)


# --------------------------------------------------------------------------
# SLIDES
# --------------------------------------------------------------------------
def s01_titre(prs):
    s = slide(prs)
    rect(s, 0, 0, 0.32, PH, fill=SLATE)
    _, tf = textbox(s, ML + 0.15, 1.95, CW, 0.4)
    para(tf, "EXPOSÉ — INTELLIGENCE ARTIFICIELLE", size=14, bold=True,
         color=SLATE, first=True, space_after=0)
    _, tf = textbox(s, ML + 0.13, 2.4, 11.0, 2.0)
    para(tf, "De ChatGPT", size=58, bold=True, color=INK, first=True,
         space_after=0, line=0.98)
    para(tf, [("au ", False, INK), ("Multi-Agent", True, SLATE)],
         size=58, bold=True, space_after=0, line=0.98)
    _, tf = textbox(s, ML + 0.15, 4.6, 10.5, 0.6)
    para(tf, "Comment l'IA passe d'un seul cerveau… à une équipe.",
         size=20, color=MUTE, first=True, space_after=0)
    labels = [("Planifier", PLAN), ("Coder", DEV), ("Tester", TEST), ("Corriger", FIX)]
    x = ML + 0.15
    for txt, col in labels:
        shape(s, MSO_SHAPE.OVAL, x, 5.7, 0.2, 0.2, fill=col)
        _, tf = textbox(s, x + 0.28, 5.66, 1.5, 0.3)
        para(tf, txt, size=12, bold=True, color=INK, first=True, space_after=0)
        x += 1.85
    return s


def s02_accroche(prs):
    s = slide(prs)
    header(s, "Pour commencer", "Une question toute simple")
    b = card(s, ML, 2.5, 8.6, 2.2, fill=tint(SLATE, 0.9), line_color=SLATE,
             line_w=1.5, radius=0.08)
    tf = b.text_frame
    para(tf, "« Qui a déjà utilisé ChatGPT,", size=26, bold=True, color=INK,
         first=True, space_after=2, align=PP_ALIGN.CENTER)
    para(tf, "Claude ou Gemini ? »", size=26, bold=True, color=INK,
         space_after=0, align=PP_ALIGN.CENTER)
    tail = shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, ML + 1.2, 4.62, 0.5, 0.4,
                 fill=tint(SLATE, 0.9), line_color=SLATE, line_w=1.5)
    tail.rotation = 180
    _, tf = textbox(s, ML, 5.5, 9.0, 1.0)
    para(tf, "On part de ce que vous connaissez déjà — puis on regarde ce qui",
         size=16, color=MUTE, first=True, space_after=2)
    para(tf, "se cache vraiment derrière ces applications.", size=16, color=MUTE,
         space_after=0)
    big = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 10.3, 2.7, 2.2, 2.0,
                fill=SLATE, radius=0.12)
    para(big.text_frame, "?", size=90, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 2)
    return s


def s03_chatgpt(prs):
    s = slide(prs)
    header(s, "Étape 1", "Comment marche ChatGPT ?")
    _, tf = textbox(s, ML, 1.95, CW, 0.6)
    para(tf, [("L'idée est étonnamment simple : l'IA ", False, INK),
              ("devine le mot suivant", True, SLATE),
              (", encore et encore, jusqu'à former une phrase.", False, INK)],
         size=17, first=True, space_after=0)
    _, tf = textbox(s, ML, 2.5, CW, 0.4)
    para(tf, "C'est comme le clavier de ton téléphone qui propose le mot d'après — mais en BEAUCOUP plus malin.",
         size=14, color=MUTE, first=True, space_after=0)
    y, h = 3.35, 1.0
    start = card(s, ML, y, 2.5, h, fill=tint(SLATE, 0.88), line_color=SLATE)
    para(start.text_frame, "Tu écris", size=15, bold=True, color=SLATE, first=True,
         space_after=2, align=PP_ALIGN.CENTER)
    para(start.text_frame, "« Le chat est… »", size=12, color=INK, space_after=0,
         align=PP_ALIGN.CENTER)
    words = ["sur", "le", "toit"]
    x = ML + 2.5 + 0.55
    wbox = 1.55
    prev_right = ML + 2.5
    line(s, prev_right + 0.04, y + h / 2, x - 0.06, y + h / 2, color=SLATE, w=2.0)
    for w in words:
        c = card(s, x, y + 0.12, wbox, h - 0.24, fill=PAPER, line_color=SLATE)
        para(c.text_frame, w, size=16, bold=True, color=INK, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        nxt = x + wbox + 0.45
        line(s, x + wbox + 0.04, y + h / 2, nxt - 0.06, y + h / 2, color=SLATE, w=2.0)
        x = nxt
    end = card(s, x, y, 2.6, h, fill=SLATE)
    para(end.text_frame, "Réponse complète", size=15, bold=True, color=PAPER,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    _, tf = textbox(s, ML + 2.5 + 0.4, y + h + 0.12, 6.5, 0.3)
    para(tf, "devine… devine… devine…", size=12, color=MUTE, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    band = rect(s, ML, 5.55, CW, 0.85, fill=tint(SLATE, 0.92))
    rect(s, ML, 5.55, 0.12, 0.85, fill=SLATE)
    para(band.text_frame,
         [("À RETENIR :   ", True, SLATE),
          ("1 question", True, INK), ("   →   ", False, MUTE),
          ("1 IA", True, INK), ("   →   ", False, MUTE),
          ("1 réponse", True, INK)],
         size=18, first=True, space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 3)
    return s


def s04_interface(prs):
    s = slide(prs)
    header(s, "Étape 2", "L'application n'est pas le cerveau")
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, "Attention au piège : ChatGPT, ce n'est pas l'intelligence elle-même. C'est juste la fenêtre où tu écris.",
         size=17, first=True, space_after=0)
    cy, ch = 2.95, 2.7
    cw = 5.2
    left = card(s, ML, cy, cw, ch, fill=PAPER, line_color=SLATE, line_w=1.5)
    tf = left.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, "L'APPLICATION", size=13, bold=True, color=SLATE, first=True, space_after=4)
    para(tf, "ChatGPT", size=30, bold=True, color=INK, space_after=8)
    para(tf, "C'est l'écran, les boutons, la zone de texte.", size=15, color=INK, space_after=3)
    para(tf, "C'est le messager qui transmet ta question.", size=15, color=INK, space_after=0)
    right = card(s, ML + cw + 1.4, cy, cw, ch, fill=SLATE, line_color=SLATE, line_w=1.5)
    tf = right.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, "LE CERVEAU", size=13, bold=True, color=tint(SLATE, 0.7), first=True, space_after=4)
    para(tf, "Le modèle (ex. GPT)", size=27, bold=True, color=PAPER, space_after=8)
    para(tf, "C'est lui qui a « appris » sur des millions de textes.", size=15, color=PAPER, space_after=3)
    para(tf, "C'est lui qui réfléchit et rédige la réponse.", size=15, color=PAPER, space_after=0)
    line(s, ML + cw + 0.1, cy + ch / 2, ML + cw + 1.3, cy + ch / 2, color=INK, w=2.4)
    _, tf = textbox(s, ML + cw + 0.0, cy + ch / 2 - 0.45, 1.4, 0.4)
    para(tf, "transmet", size=11, color=MUTE, first=True, space_after=0, align=PP_ALIGN.CENTER)
    _, tf = textbox(s, ML, cy + ch + 0.2, CW, 0.5)
    para(tf, [("L'analogie :  ", True, SLATE),
              ("le téléphone (l'application) n'est pas la personne qui répond (le cerveau).", False, INK)],
         size=16, first=True, space_after=0)
    footer(s, 4)
    return s


def s05_modeles(prs):
    s = slide(prs)
    header(s, "Étape 3", "Il existe plusieurs cerveaux")
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, "Plusieurs entreprises ont créé leur propre modèle. Chacun a ses forces.",
         size=17, first=True, space_after=0)
    data = [
        ("OpenAI", "GPT", "Polyvalent et rapide", PLAN),
        ("Anthropic", "Claude", "Doué pour les longs textes et le code", DEV),
        ("Google", "Gemini", "Comprend aussi images et vidéos", TEST),
        ("Meta", "Llama", "Gratuit et ouvert à tous", FIX),
    ]
    n = len(data)
    gap = 0.4
    cw = (CW - gap * (n - 1)) / n
    x = ML
    y, h = 2.95, 3.1
    for company, model, strength, col in data:
        c = card(s, x, y, cw, h, fill=PAPER, line_color=HAIR)
        rect(s, x, y, cw, 0.16, fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.28)
        para(tf, company.upper(), size=12, bold=True, color=col, first=True, space_after=3)
        para(tf, model, size=26, bold=True, color=INK, space_after=10)
        para(tf, strength, size=14, color=INK, space_after=0, line=1.1)
        x += cw + gap
    _, tf = textbox(s, ML, y + h + 0.2, CW, 0.5)
    para(tf, [("À retenir :  ", True, SLATE),
              ("différents cerveaux = différentes forces et faiblesses.", False, INK)],
         size=16, first=True, space_after=0)
    footer(s, 5)
    return s


def s06_transition(prs):
    s = slide(prs)
    header(s, "Le déclic", "Et si on en mettait plusieurs ?")
    card(s, ML, 2.4, 5.2, 3.6, fill=PAPER, line_color=HAIR)
    _, tf = textbox(s, ML + 0.3, 2.6, 4.6, 0.4)
    para(tf, "AUJOURD'HUI", size=13, bold=True, color=MUTE, first=True, space_after=0)
    one = shape(s, MSO_SHAPE.OVAL, ML + 0.5, 3.6, 1.3, 1.3, fill=SLATE)
    para(one.text_frame, "1 IA", size=17, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    line(s, ML + 1.95, 4.25, ML + 3.1, 4.25, color=INK, w=2.4)
    res = card(s, ML + 3.2, 3.85, 1.6, 0.8, fill=tint(SLATE, 0.9), line_color=SLATE)
    para(res.text_frame, "1 réponse", size=13, bold=True, color=INK, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    card(s, ML + 6.0, 2.4, 5.83, 3.6, fill=tint(SLATE, 0.94), line_color=SLATE, line_w=1.5)
    _, tf = textbox(s, ML + 6.3, 2.6, 5.0, 0.4)
    para(tf, "ET SI… DEMAIN", size=13, bold=True, color=SLATE, first=True, space_after=0)
    roles = [(PLAN, 3.15), (DEV, 3.95), (TEST, 4.75)]
    merge_x = ML + 9.4
    for col, cy in roles:
        o = shape(s, MSO_SHAPE.OVAL, ML + 6.4, cy, 0.7, 0.7, fill=col)
        para(o.text_frame, "IA", size=12, bold=True, color=PAPER, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        line(s, ML + 7.2, cy + 0.35, merge_x - 0.05, 4.25, color=col, w=2.0)
    out = card(s, merge_x, 3.85, 1.9, 0.85, fill=SLATE)
    para(out.text_frame, "1 réponse meilleure", size=13, bold=True, color=PAPER,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    _, tf = textbox(s, ML, 6.25, CW, 0.5)
    para(tf, "Plusieurs cerveaux qui travaillent ENSEMBLE sur le même problème : c'est le multi-agent.",
         size=16, bold=True, color=INK, first=True, space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 6)
    return s


def s07_analogie(prs):
    s = slide(prs)
    header(s, "L'analogie", "Comme un exposé en groupe")
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, "Tu connais déjà le travail d'équipe en classe. Le multi-agent, c'est exactement pareil.",
         size=17, first=True, space_after=0)
    rows = [
        ("Le chef de projet", "qui organise", "L'agent coordinateur", PLAN),
        ("Celui qui rédige", "qui produit le travail", "L'agent qui code / écrit", DEV),
        ("Celui qui relit", "qui vérifie", "L'agent vérificateur", TEST),
        ("Celui qui corrige", "qui répare les fautes", "L'agent correcteur", FIX),
    ]
    y = 2.7
    rh = 0.74
    gap = 0.1
    colL = 4.7
    colR = 4.7
    mid = 1.6
    xL = ML
    xR = ML + colL + mid
    hL = rect(s, xL, y, colL, 0.5, fill=tint(SLATE, 0.88))
    para(hL.text_frame, "DANS UN EXPOSÉ DE GROUPE", size=12, bold=True, color=SLATE,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    hR = rect(s, xR, y, colR, 0.5, fill=SLATE)
    para(hR.text_frame, "DANS UNE ÉQUIPE D'IA", size=12, bold=True, color=PAPER,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    y += 0.5 + gap
    for human, role, agent, col in rows:
        cL = card(s, xL, y, colL, rh, fill=PAPER, line_color=HAIR)
        para(cL.text_frame, [(human, True, INK), ("  —  " + role, False, MUTE)],
             size=14, first=True, space_after=0)
        line(s, xL + colL + 0.12, y + rh / 2, xR - 0.12, y + rh / 2, color=col, w=2.2)
        cR = card(s, xR, y, colR, rh, fill=tint(col, 0.86), line_color=col)
        para(cR.text_frame, agent, size=15, bold=True, color=col, first=True,
             space_after=0)
        y += rh + gap
    footer(s, 7)
    return s


def s08_agent(prs):
    s = slide(prs)
    header(s, "Définition", "C'est quoi un « agent » ?")
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, "Un agent, c'est une petite recette en trois ingrédients :",
         size=17, first=True, space_after=0)
    y, h = 3.0, 2.0
    items = [
        ("UN CERVEAU", "un modèle (GPT, Claude…) qui réfléchit", SLATE),
        ("UN RÔLE", "une mission précise (« tu testes le code »)", DEV),
        ("DES OUTILS", "de quoi agir (chercher, lancer un test…)", TEST),
    ]
    cw = 3.15
    xs = [ML, ML + 3.55, ML + 7.1]
    for (title, desc, col), x in zip(items, xs):
        c = card(s, x, y, cw, h, fill=PAPER, line_color=col, line_w=1.5)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        para(tf, title, size=15, bold=True, color=col, first=True, space_after=6)
        para(tf, desc, size=14, color=INK, space_after=0, line=1.12)
    for px in [ML + 3.18, ML + 6.73]:
        _, tf = textbox(s, px, y + h / 2 - 0.35, 0.4, 0.7, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "+", size=40, bold=True, color=MUTE, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
    res = card(s, ML + 10.45, y, 1.38, h, fill=SLATE)
    para(res.text_frame, "= un AGENT", size=20, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    _, tf = textbox(s, ML, y + h + 0.35, CW, 0.5)
    para(tf, [("En clair :  ", True, SLATE),
              ("un cerveau à qui on donne un rôle et des outils devient un travailleur spécialisé.", False, INK)],
         size=16, first=True, space_after=0)
    footer(s, 8)
    return s


def s09_coulisses(prs):
    s = slide(prs)
    header(s, "En coulisses", "« Fais-moi un jeu » : qui fait quoi ?")
    cmd = card(s, ML, 1.95, CW, 0.6, fill=tint(SLATE, 0.9), line_color=SLATE)
    para(cmd.text_frame, [("TA DEMANDE :   ", True, SLATE),
                          ("« Fais-moi un jeu de Tetris »", False, INK)],
         size=15, first=True, space_after=0, align=PP_ALIGN.CENTER)
    agents = [
        ("1", "Le Planificateur", "« Je découpe le travail en étapes claires. »", PLAN),
        ("2", "Le Développeur", "« J'écris le code du jeu. »", DEV),
        ("3", "Le Testeur", "« Je vérifie que le jeu marche vraiment. »", TEST),
        ("4", "Le Correcteur", "« Je répare les erreurs trouvées. »", FIX),
    ]
    y = 2.75
    rh = 0.78
    gap = 0.18
    for num, name, quote, col in agents:
        card(s, ML, y, CW, rh, fill=tint(col, 0.9), line_color=col)
        badge = shape(s, MSO_SHAPE.OVAL, ML + 0.18, y + rh / 2 - 0.26, 0.52, 0.52, fill=col)
        para(badge.text_frame, num, size=20, bold=True, color=PAPER, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        _, tf = textbox(s, ML + 0.95, y + 0.1, 3.3, rh - 0.2, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, name, size=16, bold=True, color=col, first=True, space_after=0)
        _, tf = textbox(s, ML + 4.4, y + 0.1, CW - 4.6, rh - 0.2, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, quote, size=15, color=INK, first=True, space_after=0)
        if num != "4":
            line(s, ML + 0.44, y + rh, ML + 0.44, y + rh + gap, color=MUTE, w=2.0)
        y += rh + gap
    _, tf = textbox(s, ML, y + 0.02, CW, 0.4)
    para(tf, "Chacun fait sa partie, puis passe le relais : c'est une vraie équipe.",
         size=15, bold=True, color=INK, first=True, space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 9)
    return s


def s10_framework(prs):
    s = slide(prs)
    header(s, "L'organisation", "Le « chef d'orchestre »")
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, [("Un ", False, INK), ("framework", True, SLATE),
              (", c'est ce qui décide qui travaille, dans quel ordre, et qui parle à qui.", False, INK)],
         size=17, first=True, space_after=0)
    pw = 3.7
    gap = 0.36
    py = 2.75
    ph = 3.5
    xs = [ML, ML + pw + gap, ML + 2 * (pw + gap)]
    titles = ["À LA CHAÎNE", "EN PARALLÈLE", "AVEC UN CHEF"]
    subs = ["l'un après l'autre", "tous en même temps", "un superviseur répartit"]
    cols = [PLAN, DEV, FIX]
    for x, t, sub, col in zip(xs, titles, subs, cols):
        card(s, x, py, pw, ph, fill=PAPER, line_color=HAIR)
        rect(s, x, py, pw, 0.14, fill=col)
        _, tf = textbox(s, x + 0.2, py + 0.3, pw - 0.4, 0.5)
        para(tf, t, size=14, bold=True, color=col, first=True, space_after=1)
        para(tf, sub, size=11, color=MUTE, space_after=0)

    cy = py + 1.5
    bx = xs[0] + 0.45
    for i, lab in enumerate(["A", "B", "C"]):
        o = shape(s, MSO_SHAPE.OVAL, bx, cy, 0.6, 0.6, fill=PLAN)
        para(o.text_frame, lab, size=15, bold=True, color=PAPER, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        if i < 2:
            line(s, bx + 0.62, cy + 0.3, bx + 0.92, cy + 0.3, color=INK, w=2.0)
        bx += 0.95
    _, tf = textbox(s, xs[0] + 0.2, cy + 0.95, pw - 0.4, 0.4)
    para(tf, "Résultat à la fin", size=11, color=INK, first=True, space_after=0,
         align=PP_ALIGN.CENTER)

    px = xs[1] + 0.35
    ys = [cy - 0.45, cy + 0.15, cy + 0.75]
    syn_x = xs[1] + pw - 1.15
    for lab, yy in zip(["A", "B", "C"], ys):
        o = shape(s, MSO_SHAPE.OVAL, px, yy, 0.55, 0.55, fill=DEV)
        para(o.text_frame, lab, size=14, bold=True, color=PAPER, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        line(s, px + 0.57, yy + 0.27, syn_x - 0.04, cy + 0.3, color=DEV, w=1.8)
    syn = card(s, syn_x, cy + 0.02, 0.95, 0.6, fill=DEV)
    para(syn.text_frame, "synthèse", size=11, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)

    boss_cx = xs[2] + pw / 2
    boss = shape(s, MSO_SHAPE.OVAL, boss_cx - 0.6, cy - 0.5, 1.2, 0.6, fill=FIX)
    para(boss.text_frame, "chef", size=12, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    kid_y = cy + 0.65
    kxs = [xs[2] + 0.55, xs[2] + pw / 2 - 0.27, xs[2] + pw - 1.1]
    for kx in kxs:
        shape(s, MSO_SHAPE.OVAL, kx, kid_y, 0.54, 0.54, fill=tint(FIX, 0.25))
        line(s, boss_cx, cy + 0.12, kx + 0.27, kid_y - 0.02, color=FIX, w=1.8, arrow=False)
    footer(s, 10)
    return s


def s11_review(prs):
    s = slide(prs)
    header(s, "Exemple concret", "Faire relire du code par une équipe", kcolor=TEST)
    _, tf = textbox(s, ML, 1.95, CW, 0.5)
    para(tf, "Avant de valider un programme, on le fait passer par plusieurs vérifications successives.",
         size=17, first=True, space_after=0)
    cx = ML + 1.3
    top = card(s, cx, 2.55, 4.2, 0.6, fill=tint(SLATE, 0.9), line_color=SLATE)
    para(top.text_frame, "Le code à vérifier", size=15, bold=True, color=INK,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    stations = [
        ("Agent Sécurité", "« Aucune faille trouvée. »", TEST),
        ("Agent Style", "« C'est bien rangé et lisible. »", PLAN),
        ("Agent Tests", "« Tout fonctionne. »", DEV),
    ]
    y = 3.15
    sh_h = 0.6
    gap = 0.24
    for name, verdict, col in stations:
        line(s, cx + 2.1, y, cx + 2.1, y + gap, color=MUTE, w=2.0)
        y += gap
        card(s, cx, y, 4.2, sh_h, fill=tint(col, 0.88), line_color=col)
        check_badge(s, cx + 0.42, y + sh_h / 2, r=0.17, color=col)
        _, tf = textbox(s, cx + 0.78, y + 0.04, 3.3, sh_h - 0.08, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, name, size=14, bold=True, color=col, first=True, space_after=0)
        y += sh_h
    line(s, cx + 2.1, y, cx + 2.1, y + gap, color=MUTE, w=2.0)
    y += gap
    final = card(s, cx, y, 4.2, 0.66, fill=TEST)
    para(final.text_frame, "Code validé", size=18, bold=True, color=PAPER, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    rx = ML + 6.4
    _, tf = textbox(s, rx, 2.7, CW - 6.4 - 0.3, 0.4)
    para(tf, "CE QUE CHAQUE AGENT DIT", size=12, bold=True, color=MUTE, first=True, space_after=0)
    yy = 3.25
    for name, verdict, col in stations:
        c = card(s, rx, yy, CW - 6.4, 0.9, fill=PAPER, line_color=HAIR)
        rect(s, rx, yy, 0.12, 0.9, fill=col)
        tf = c.text_frame
        para(tf, name, size=13, bold=True, color=col, first=True, space_after=2)
        para(tf, verdict, size=14, color=INK, space_after=0)
        yy += 1.04
    footer(s, 11)
    return s


def s12_comparaison(prs):
    s = slide(prs)
    header(s, "Comparons", "ChatGPT seul ou une équipe ?")
    rows = [
        ("Architecture", "1 seul cerveau", "Plusieurs cerveaux"),
        ("Façon de faire", "Réponse directe", "Collaboration"),
        ("Tâches", "Simples, rapides", "Compliquées, en plusieurs étapes"),
        ("Coût", "Un seul appel", "Plusieurs appels (donc plus cher)"),
    ]
    x0 = ML
    wlab = 2.6
    wcol = (CW - wlab) / 2
    xA = x0 + wlab
    xB = xA + wcol
    y = 2.35
    hh = 0.62
    rect(s, x0, y, wlab, hh, fill=CREAM)
    a = rect(s, xA, y, wcol, hh, fill=SLATE)
    para(a.text_frame, "ChatGPT tout seul", size=15, bold=True, color=PAPER,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    b = rect(s, xB, y, wcol, hh, fill=DEV)
    para(b.text_frame, "Une équipe d'agents", size=15, bold=True, color=PAPER,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    y += hh
    rh = 0.82
    for i, (lab, va, vb) in enumerate(rows):
        bg = PAPER if i % 2 == 0 else tint(SLATE, 0.95)
        lc = rect(s, x0, y, wlab, rh, fill=tint(SLATE, 0.9))
        para(lc.text_frame, lab, size=14, bold=True, color=SLATE, first=True,
             space_after=0)
        ca = rect(s, xA, y, wcol, rh, fill=bg, line_color=HAIR, line_w=0.75)
        para(ca.text_frame, va, size=14, color=INK, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
        cb = rect(s, xB, y, wcol, rh, fill=tint(DEV, 0.92) if i % 2 else tint(DEV, 0.86),
                  line_color=HAIR, line_w=0.75)
        para(cb.text_frame, vb, size=14, color=INK, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
        y += rh
    _, tf = textbox(s, ML, y + 0.2, CW, 0.5)
    para(tf, [("En résumé :  ", True, SLATE),
              ("plus puissant pour les tâches difficiles… mais plus lent et plus coûteux.", False, INK)],
         size=16, first=True, space_after=0)
    footer(s, 12)
    return s


def s13_pyramide(prs):
    s = slide(prs)
    header(s, "À retenir", "Tout tient dans une pyramide")
    cx = ML + 3.3
    base_y = 5.6
    widths = [5.4, 3.9, 2.4]
    labels = [("MODÈLE", "le cerveau (GPT, Claude…)", SLATE),
              ("AGENT", "cerveau + rôle + outils", DEV),
              ("ÉQUIPE D'AGENTS", "le système multi-agent", TEST)]
    h = 1.05
    gap = 0.12
    for i, (w, (lab, desc, col)) in enumerate(zip(widths, labels)):
        y = base_y - i * (h + gap)
        x = cx + (widths[0] - w) / 2
        sh = shape(s, MSO_SHAPE.TRAPEZOID if i > 0 else MSO_SHAPE.RECTANGLE,
                   x, y, w, h, fill=col)
        if i > 0:
            try:
                sh.adjustments[0] = 0.25
            except Exception:
                pass
        para(sh.text_frame, lab, size=18, bold=True, color=PAPER, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        line(s, x + w + 0.05, y + h / 2, cx + widths[0] + 0.45, y + h / 2,
             color=col, w=1.6, arrow=False)
        _, tf = textbox(s, cx + widths[0] + 0.6, y + h / 2 - 0.25, 2.85, 0.7,
                        anchor=MSO_ANCHOR.MIDDLE)
        para(tf, lab, size=13, bold=True, color=col, first=True, space_after=1)
        para(tf, desc, size=12, color=INK, space_after=0)
    _, tf = textbox(s, cx, base_y + h + 0.25, widths[0], 0.5)
    para(tf, "Cerveau  →  Rôle  →  Équipe", size=18, bold=True, color=INK,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 13)
    return s


def s14_conclusion(prs):
    s = slide(prs)
    header(s, "Conclusion", "Ce qu'il faut vraiment retenir")
    band = card(s, ML, 2.05, CW, 1.25, fill=SLATE, radius=0.06)
    tf = band.text_frame
    para(tf, "L'IA, ce n'est pas juste un cerveau…", size=22, bold=True, color=PAPER,
         first=True, space_after=2, align=PP_ALIGN.CENTER)
    para(tf, [("c'est ", False, PAPER), ("comment on l'organise.", True, tint(TEST, 0.4))],
         size=22, bold=True, space_after=0, align=PP_ALIGN.CENTER)
    q = card(s, ML, 3.6, CW, 0.7, fill=tint(FIX, 0.85), line_color=FIX, line_w=1.5)
    para(q.text_frame, "Question pour vous : « Plus d'agents = toujours mieux ? »",
         size=17, bold=True, color=INK, first=True, space_after=0, align=PP_ALIGN.CENTER)
    answers = [
        ("Non", "ça coûte du temps et de l'argent.", PLAN),
        ("Non", "trop d'agents = on s'embrouille.", DEV),
        ("Oui, si…", "l'organisation est bien pensée.", TEST),
    ]
    cw = (CW - 2 * 0.4) / 3
    x = ML
    y = 4.6
    hh = 1.5
    for head, body, col in answers:
        c = card(s, x, y, cw, hh, fill=PAPER, line_color=col, line_w=1.5)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        para(tf, head, size=20, bold=True, color=col, first=True, space_after=5)
        para(tf, body, size=15, color=INK, space_after=0, line=1.12)
        x += cw + 0.4
    _, tf = textbox(s, ML, y + hh + 0.18, CW, 0.4)
    para(tf, "Ce qui compte, ce n'est pas le nombre — c'est la bonne organisation.",
         size=15, bold=True, color=SLATE, first=True, space_after=0, align=PP_ALIGN.CENTER)
    footer(s, 14)
    return s


# --------------------------------------------------------------------------
def main():
    prs = new_deck()
    for fn in (s01_titre, s02_accroche, s03_chatgpt, s04_interface, s05_modeles,
               s06_transition, s07_analogie, s08_agent, s09_coulisses,
               s10_framework, s11_review, s12_comparaison, s13_pyramide,
               s14_conclusion):
        fn(prs)
    out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "presentation_IA.pptx")
    prs.save(out)
    print("Saved", out, "-", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    main()
